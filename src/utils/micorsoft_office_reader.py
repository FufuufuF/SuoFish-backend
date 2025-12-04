from io import BytesIO
from typing import Union

from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table

from pypdf import PdfReader

from pptx import Presentation


class MicrosoftOfficeReader:
    """Office 文档解析器，支持从文件路径或字节数据读取"""
    
    @staticmethod
    def read_docx(source: Union[str, bytes, BytesIO]) -> str:
        """
        读取 Word 文档
        
        Args:
            source: 文件路径、字节数据或 BytesIO 对象
        """
        try:
            # 如果是 bytes，转换为 BytesIO
            if isinstance(source, bytes):
                source = BytesIO(source)
            
            doc = Document(source)
            full_text = ""
            for block in doc.iter_inner_content():
                if isinstance(block, Paragraph):
                    text = block.text.strip()
                    if text:
                        full_text += text + '\n'
                elif isinstance(block, Table):
                    for row in block.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                full_text += text + '\t'
                        full_text += '\n'
            return full_text
        except Exception as e:
            return f"读取 Word 文档失败: {e}"

    @staticmethod
    def read_pdf(source: Union[str, bytes, BytesIO]) -> str:
        """
        读取 PDF 文档
        
        Args:
            source: 文件路径、字节数据或 BytesIO 对象
        """
        try:
            # 如果是 bytes，转换为 BytesIO
            if isinstance(source, bytes):
                source = BytesIO(source)
            
            reader = PdfReader(source)
            full_text = ""
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    full_text += text + '\n'
            return full_text
        except Exception as e:
            return f"读取 PDF 失败: {e}"

    @staticmethod
    def read_pptx(source: Union[str, bytes, BytesIO]) -> str:
        """
        读取 PowerPoint 文档
        
        Args:
            source: 文件路径、字节数据或 BytesIO 对象
        """
        try:
            # 如果是 bytes，转换为 BytesIO
            if isinstance(source, bytes):
                source = BytesIO(source)
            
            presentation = Presentation(source)
            full_text = []
            for slide_num, slide in enumerate(presentation.slides):
                devide_str = f"\n--- 第 {slide_num + 1} 张幻灯片 ---"
                slide_content = [devide_str]
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            slide_content.append(text)
                    
                    if shape.has_table:
                        table_content = []
                        table = shape.table
                        for row in table.rows:
                            row_text = [cell.text_frame.text.strip() for cell in row.cells]
                            table_content.append(" | ".join(row_text))
                        slide_content.append("\n[表格内容]:\n" + "\n".join(table_content))

                full_text.append('\n'.join(slide_content))
            return '\n'.join(full_text)
            
        except Exception as e:
            return f"读取 PPT 失败: {e}"
